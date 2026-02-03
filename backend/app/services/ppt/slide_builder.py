import os
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE


from app.services.ppt.template_loader import load_template
from app.services.ppt.chart_renderer import render_chart
from pptx.dml.color import RGBColor


EMU_PER_INCH = 914400
PX_PER_INCH = 96


def px_to_emu(px: int) -> int:
    return int(px / PX_PER_INCH * EMU_PER_INCH)


# --------------------------------------------------
# Find branded layout (green header slide)
# --------------------------------------------------

def _find_branded_layout(prs):

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for layout in prs.slide_layouts:

        for shp in layout.shapes:

            if shp.is_placeholder:
                continue

            if shp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:

                if shp.top < slide_h * 0.18 and shp.width > slide_w * 0.6:
                    return layout

    return prs.slide_layouts[-1]


# --------------------------------------------------
# PPT BUILDER (TEMPLATE-DRIVEN)
# --------------------------------------------------

def build_ppt_from_slides(project, slides, output_path, debug_mode: bool = False):

    prs = load_template()

    base_layout = _find_branded_layout(prs)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # ---- layout constants (px) ----
    SIDE_MARGIN = px_to_emu(80)
    HEADER_Y = px_to_emu(18)
    HEADER_H = px_to_emu(56)
    CHART_TOP = px_to_emu(155)
    FOOTER_CLEARANCE = px_to_emu(80)

    for idx, slide in enumerate(slides):

        slide_json = slide.slide_json
        
        # -----------------------------
        # VALIDATION (Hard Failure)
        # -----------------------------
        if not slide_json.get("width") or not slide_json.get("height"):
             # Technically PPTX template drives this, but if the slide JSON implies
             # a specific layout intent that's missing, we should know. 
             # Actually, for elements, let's enforce validation there.
             raise ValueError("Slide JSON is missing 'width' or 'height'")

        if not slide_json.get("elements"):
            # Empty slide is suspicious but not necessarily fatal? 
            # Reqs say "slide JSON missing width/height". 
            # But the slide dimensions come from template generally.
            # I will check if elements are valid.
            raise ValueError("Slide JSON has no elements. Aborting export.")

        # Always create new slide to ensure clean state
        ppt_slide = prs.slides.add_slide(base_layout)

        # -----------------------------
        # Render elements
        # -----------------------------
        for el in slide_json.get("elements", []):

            el_type = el.get("type")

            # ---------------- TITLE ----------------
            if el_type == "title":

                x = SIDE_MARGIN
                y = HEADER_Y
                w = slide_width - SIDE_MARGIN * 2
                h = HEADER_H

                box = ppt_slide.shapes.add_textbox(x, y, w, h)
                if el.get("boxed"):
                  fill = box.fill
                  fill.solid()
                  fill.fore_color.rgb = RGBColor(235, 245, 241)

                  box.line.color.rgb = RGBColor(46, 139, 87)

                  box.line.width = Pt(1.5)

                  box.shadow.inherit = False

                tf = box.text_frame
                tf.clear()

                p = tf.paragraphs[0]
                p.text = el.get("value", "")
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER

            # ---------------- TEXT ----------------
            elif el_type == "text":

                x = px_to_emu(el.get("x", 0))
                y = px_to_emu(el.get("y", 0))
                w = px_to_emu(el.get("width", 600))
                h = px_to_emu(el.get("height", 120))

                box = ppt_slide.shapes.add_textbox(x, y, w, h)
                tf = box.text_frame
                tf.clear()

                p = tf.paragraphs[0]
                p.text = el.get("value", "")
                p.font.size = Pt(el.get("fontSize", 16))
                p.alignment = PP_ALIGN.LEFT

            # ---------------- CHART ----------------
            elif el_type == "chart":

                name = (el.get("datasetName") or "").lower()

                # ---------------- SIZE RULES ----------------

                if name.endswith("_lo"):
                    lo_count = len(el.get("preview", []))
                    BASE = 360
                    PER_ROW = 45
                    chart_w = px_to_emu(920)
                    chart_h = px_to_emu(min(900, BASE + lo_count * PER_ROW))

                elif name.endswith("_qlvl"):
                    chart_w = px_to_emu(760)
                    chart_h = px_to_emu(440)

                elif "reg_vs_part" in name:
                    chart_w = px_to_emu(700)
                    chart_h = px_to_emu(420)

                elif "summary" in name:
                    chart_w = px_to_emu(760)
                    chart_h = px_to_emu(430)

                else:
                    chart_w = px_to_emu(760)
                    chart_h = px_to_emu(430)

                # ---------------- CENTER + CLAMP ----------------

                x = int((slide_width - chart_w) / 2)

                max_h = slide_height - CHART_TOP - FOOTER_CLEARANCE
                if chart_h > max_h:
                    chart_h = max_h

                y = CHART_TOP

                # inject geometry back so renderer remains local
                el["x"] = int(x / EMU_PER_INCH * PX_PER_INCH)
                el["y"] = int(y / EMU_PER_INCH * PX_PER_INCH)
                el["width"] = int(chart_w / EMU_PER_INCH * PX_PER_INCH)
                el["height"] = int(chart_h / EMU_PER_INCH * PX_PER_INCH)

                # render chart
                
                # Header flag OR env var
                use_debug = debug_mode or (os.environ.get("EXPORT_DEBUG") == "true")

                render_chart(
                    ppt_slide=ppt_slide,
                    element=el,
                    debug_mode=use_debug,
                )

                # DEBUG: Bounding Box
                if use_debug:
                    debug_box = ppt_slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        px_to_emu(el["x"]),
                        px_to_emu(el["y"]),
                        px_to_emu(el["width"]),
                        px_to_emu(el["height"]),
                    )

                    debug_box.fill.background()
                    debug_box.line.color.rgb = RGBColor(0, 0, 255)
                    debug_box.line.width = Pt(2)
                    # No text in box
                    debug_box.text_frame.clear()

        # -----------------------------
        # Footer (Version Stamp)
        # -----------------------------
        _add_version_footer(ppt_slide, slide_width, slide_height)

    prs.save(output_path)
    return output_path


def _add_version_footer(slide, slide_width, slide_height):
    from datetime import datetime
    
    version = "1.2"
    date_str = datetime.now().strftime("%Y-%m-%d")
    text = f"Generated by Dashboard Engine v{version} | {date_str}"

    # Positioning: Bottom Right
    # Using small font, subdued color
    
    margin = px_to_emu(20)
    w = px_to_emu(400)
    h = px_to_emu(30)
    x = slide_width - w - margin
    y = slide_height - h - margin

    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT
