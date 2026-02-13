print("USING CHART RENDERER FROM:", __file__)



from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
from pptx.chart.data import ChartData
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

import re
import math
import os


EMU_PER_INCH = 914400
PX_PER_INCH = 96


def px_to_emu(px: int) -> int:
    return int(px / PX_PER_INCH * EMU_PER_INCH)


# --------------------------------------------------
# Chart mapping
# --------------------------------------------------

CHART_MAP = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "grouped_bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar-horizontal": XL_CHART_TYPE.BAR_CLUSTERED,
    "pie": XL_CHART_TYPE.PIE,
}

DEFAULT_BAR_COLOR = RGBColor(79, 129, 189)


# --------------------------------------------------
# Dataset type detection
# --------------------------------------------------

def detect_dataset_type(element):
    """
    Returns the classification for the dataset.
    Prioritizes explicit 'datasetFamily' from the element metadata.
    """
    family = element.get("datasetFamily")
    
    # Do not trust generic family blindly
    if family and family != "generic":
        return family


    name = (element.get("datasetName") or "").lower()

    if "sub_wise" in name or "subwise" in name:
        return "subwise"

    if name.endswith("_lo") or "learning_outcome" in name:
        return "lo"

    if name.endswith("_qlvl") or "difficulty" in name:
        return "qlvl"

    if "reg_vs_part_grade" in name or ("reg" in name and "part" in name and "grade" in name):
        return "reg_vs_part_grade"

    if "reg_vs_part_schl" in name or "reg_vs_part_school" in name:
        return "reg_vs_part_school"

    if "perf_summary" in name or "summary" in name:
        return "perf_summary"

    return "generic"


# --------------------------------------------------
# Percent column detection (NEW)
# --------------------------------------------------

def _detect_numeric_percent_column(columns, preview):
    """
    Fallback: find numeric column that looks like percent.
    """
    for c in columns:
        vals = []
        for r in preview:
            try:
                # Better 0-1 decimal handling
                raw = str(r.get(c))
                if "." in raw and float(raw) <= 1:
                    v = int(round(float(raw) * 100))
                else:
                    v = normalize_percent(raw)
                vals.append(v)
            except Exception:
                continue

        if not vals:
            continue

        mn = min(vals)
        mx = max(vals)

        # Only treat as percent if the column name clearly indicates it
        if 0 <= mn and mx <= 100:
            name = c.lower()
            if "%" in name or "percent" in name or "percentage" in name:
                return c

    return None


def normalize_percent(val):
    """
    Centralized percent handling. Strips symbols, rounds to nearest integer.
    """
    try:
        # Handle None or empty strings
        if val is None or str(val).strip() == "":
            return 0
        # Strip % and common numeric decorators
        clean_val = str(val).replace("%", "").replace("(", "").replace(")", "").strip()
        return int(round(float(clean_val)))
    except (ValueError, TypeError):
        return 0


def _safe_count(val):
    """
    Safe numeric parser for counts. Rounds and strips junk characters.
    """
    try:
        if val is None:
            return 0
        return int(round(float(str(val).replace("%", "").strip())))
    except Exception:
        return 0


# --------------------------------------------------
# LO label cleaner
# --------------------------------------------------

LO_PREFIXES = [
    r"the student will be able to",
    r"students will be able to",
    r"student will be able to",
    r"learner will be able to",
]


def clean_lo_label(text: str, max_words: int = 6) -> str:

    t = str(text).lower().strip()

    for p in LO_PREFIXES:
        t = re.sub(p, "", t).strip()

    t = t.strip(" .:")

    words = t.split()

    if len(words) > max_words:
        words = words[:max_words]

    return " ".join(words).title()


# --------------------------------------------------
# MAIN ENTRY
# --------------------------------------------------

def render_chart(ppt_slide, element, debug_mode: bool = False):

    x = px_to_emu(element.get("x", 0))
    y = px_to_emu(element.get("y", 0))
    w = px_to_emu(element.get("width", 500))
    h = px_to_emu(element.get("height", 350))

    chart_type = element.get("chartType")
    preview = element.get("preview") or []

    if not preview:
        raise ValueError(
            f"Chart Error: Dataset '{element.get('datasetName')}' has empty preview data. "
            "Cannot render chart."
        )

    dataset_type = detect_dataset_type(element)
    metric_types = element.get("metricTypes") or {}

    print(
        f"[RENDERING] Dataset: {element.get('datasetName')} | "
        f"Type: {dataset_type} | MetricTypes: {metric_types}"
    )

    columns = list(preview[0].keys())
    label_col = columns[0]

    # --------------------------------------------------
    # Percent column resolution (FIXED)
    # --------------------------------------------------

    # For reg_vs_part_grade, ONLY use schema-classified percent columns
    if dataset_type == "reg_vs_part_grade":
        percent_col = next(
            (c for c in columns if metric_types.get(c) == "percent"),
            None,
        )
    else:
        percent_col = next(
            (c for c in columns if metric_types.get(c) == "percent"),
            None,
        )

        if not percent_col:
            percent_col = _detect_numeric_percent_column(columns, preview)

        if not percent_col:
            percent_col = next(
                (c for c in columns if "%" in c or "percent" in c.lower()),
                None,
            )

    # --------------------------------------------------
    # VALIDATION
    # --------------------------------------------------

    PCT_REQUIRED_TYPES = {
        "lo",
        "qlvl",
        "perf_summary",
        "subwise",
    }

    if dataset_type in PCT_REQUIRED_TYPES and not percent_col:
        raise ValueError(
            f"Chart Error: Dataset '{element.get('datasetName')}' (type: {dataset_type}) "
            "requires a percentage column but none was detected. "
            f"Found columns: {columns}"
        )

    # --------------------------------------------------
    # Metric selection
    # --------------------------------------------------

    if dataset_type == "reg_vs_part_grade":
        # Only include count columns as bars. ignore Participation %
        metric_cols = [
            c for c in columns
            if metric_types.get(c) in ("registered", "participated")
        ]
        print(f"[DEBUG] reg_vs_part_grade DETECTED. metric_cols={metric_cols}")

    elif dataset_type in ("lo", "qlvl", "perf_summary", "subwise"):
        metric_cols = [percent_col] if percent_col else []
        if dataset_type == "subwise":
            print(f"[DEBUG] subwise DETECTED. metric_cols={metric_cols}")
        if not metric_cols and columns[1:]:
             # Last resort for percent types, but avoid Participation % if possible
             metric_cols = [columns[1]]

    else:
        metric_cols = columns[1:]

    print(f"DATASET TYPE = {dataset_type}")
    print(f"METRIC_COLS = {metric_cols}")

    # Remove dynamic fallback for reg_vs_part_grade to ensure only counts are plotted
    if not metric_cols and dataset_type != "reg_vs_part_grade":
        metric_cols = columns[1:]

    print(f"[PRE-GUARD] dataset_type={dataset_type}, metric_cols={metric_cols}, percent_col={percent_col}")

    # FINAL GUARD: reg_vs_part charts must NEVER plot percent series
    if dataset_type == "reg_vs_part_grade" and percent_col in metric_cols:
        metric_cols = [c for c in metric_cols if c != percent_col]
        print(f"[GUARD] Removed percent column '{percent_col}' from metric_cols. Final: {metric_cols}")



    chart_data = ChartData()

    # ==================================================
    # PIE
    # ==================================================

    if chart_type == "pie":
        chart = _build_pie_chart(
            ppt_slide,
            x,
            y,
            w,
            h,
            chart_data,
            preview,
            label_col,
            metric_cols,
        )

        _format_pie(chart)
        return

    # ==================================================
    # BAR / COLUMN BASE
    # ==================================================

    categories = []

    for r in preview:
        label = str(r.get(label_col))
        if dataset_type == "lo":
            label = clean_lo_label(label)
        categories.append(label)

    chart_data.categories = categories

    for col in metric_cols:
        vals = []
        is_pct = (col == percent_col) or (metric_types.get(col) == "percent")
        
        if dataset_type == "subwise":
            print("[SUBWISE] forcing percent-only plotting")

        for r in preview:
            raw_val = r.get(col)
            # Force normalization for subwise OR if identified as percent
            if dataset_type == "subwise" or is_pct:
                vals.append(normalize_percent(raw_val))
            elif dataset_type == "reg_vs_part_grade":
                vals.append(_safe_count(raw_val))
            else:
                try:
                    vals.append(float(str(raw_val).replace("%", "")))
                except Exception:
                    vals.append(0)
        chart_data.add_series(col, vals)

    pptx_chart_type = CHART_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart = ppt_slide.shapes.add_chart(
        pptx_chart_type,
        x,
        y,
        w,
        h,
        chart_data,
    ).chart

    print(f"SERIES IN CHART = {[s.name for s in chart.series]}")

    chart.has_title = False

    plot = chart.plots[0]

    # FORCE percent axis formatting when metric is percent
    if dataset_type in ("perf_summary", "summary", "lo", "qlvl"):
        val_axis = chart.value_axis
        val_axis.minimum_scale = 0.0
        val_axis.maximum_scale = 100.0
        val_axis.major_unit = 20.0

        plot.has_data_labels = True
        labels = plot.data_labels
        labels.number_format = '0"%"'

    # --------------------------------------------------
    # Dispatch formatting
    # --------------------------------------------------

    if dataset_type == "lo":
        _format_lo(chart, plot)

    elif dataset_type == "qlvl":
        _format_qlvl(chart, plot)

    elif dataset_type == "reg_vs_part_grade":
        _format_reg_vs_part_grade(chart, plot, preview, percent_col)

    elif dataset_type == "perf_summary":
        _format_perf_summary(chart, plot)

    elif dataset_type == "subwise":
        _format_subwise(chart, plot)

    elif any(metric_types.get(c) == "percent" for c in metric_cols):
        _format_perf_summary(chart, plot)
    else:
        _format_generic(chart, plot)

    _apply_bar_colors(chart)

    if debug_mode or os.environ.get("EXPORT_DEBUG") == "true":
        _render_debug_overlay(ppt_slide, x, y, w, h, element, dataset_type)


# --------------------------------------------------
# PIE BUILDER
# --------------------------------------------------

def _build_pie_chart(
    ppt_slide,
    x,
    y,
    w,
    h,
    chart_data,
    preview,
    label_col,
    metric_cols,
):

    cats = []
    vals = []

    metric = metric_cols[0]

    for r in preview:
        cats.append(str(r.get(label_col)))
        # ALWAYS use normalize_percent for PIE as they are essentially distributions/percentages
        vals.append(normalize_percent(r.get(metric)))

    chart_data.categories = cats
    chart_data.add_series(metric, vals)

    chart = ppt_slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        x,
        y,
        w,
        h,
        chart_data,
    ).chart

    return chart


# --------------------------------------------------
# FORMATTERS
# --------------------------------------------------

def _format_pie(chart):

    plot = chart.plots[0]
    plot.has_data_labels = True

    labels = plot.data_labels
    labels.show_category_name = True
    labels.show_value = True
    labels.show_percentage = True
    labels.font.size = Pt(11)
    labels.font.bold = False
    labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    chart.has_legend = False


# ---------------- GRADE WISE -----------------------

def _format_reg_vs_part_grade(chart, plot, preview, percent_col):

    # sensible spacing
    plot.gap_width = 220
    plot.overlap = -30

    cat_axis = chart.category_axis
    val_axis = chart.value_axis

    cat_axis.tick_labels.font.size = Pt(10)
    val_axis.tick_labels.font.size = Pt(10)

    val_axis.minimum_scale = 0.0

    # Use only count series (Registered, Participated) for axis scaling, NEVER Participation %
    series_names = [s.name for s in chart.series if s.name != percent_col]
    if not series_names:
        series_names = [s.name for s in chart.series]

    # STICK TO COUNTS: Ignore percent columns entirely for axis calculation
    max_val = max(
        max(_safe_count(r.get(col)) for r in preview)
        for col in series_names
    )

    val_axis.maximum_scale = float(math.ceil(max_val / 50) * 50)
    val_axis.major_unit = 50.0

    # Remove global plot-level label settings to handle per-series
    plot.has_data_labels = False

    for series in chart.series:
        
        # Enable labels for this series
        series.has_data_labels = True
        labels = series.data_labels
        labels.font.size = Pt(9)
        labels.font.bold = False

        # Stagger positions based on series type to avoid collision
        if series.name.lower().startswith("participated"):
            labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
            labels.font.color.rgb = RGBColor(0, 0, 0)
        else:
            labels.position = XL_DATA_LABEL_POSITION.INSIDE_END
            labels.font.color.rgb = RGBColor(0, 0, 0)

        for p_idx, point in enumerate(series.points):

            dl = point.data_label
            
            tf = dl.text_frame
            tf.clear()

            p = tf.paragraphs[0]
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)

            value = series.values[p_idx]
            raw = preview[p_idx].get(percent_col)

            # FORCE normalized integer % for reg_vs_part labels
            percent = normalize_percent(raw)

            if series.name.lower().startswith("participated"):
                # Ensure no decimals in the label text - format as count(percentage%)
                p.text = f"{int(round(value))}({percent}%)"
            else:
                p.text = f"{int(round(value))}"

    _set_axis_titles(chart, "Grade", "Students")

# ---------------- LO ------------------------------

def _format_lo(chart, plot):

    plot.gap_width = 150
    plot.overlap = -10

    cat_axis = chart.category_axis
    val_axis = chart.value_axis

    cat_axis.tick_labels.font.size = Pt(9)
    
    val_axis.tick_labels.font.size = Pt(9)
    val_axis.minimum_scale = 0.0
    val_axis.maximum_scale = 100.0
    val_axis.major_unit = 20.0

    plot.has_data_labels = True
    labels = plot.data_labels
    labels.font.size = Pt(9)
    labels.number_format = '0"%"'
    labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    _set_axis_titles(chart, "Learning Outcome", "Percentage")



def _format_qlvl(chart, plot):

    plot.gap_width = 150

    cat_axis = chart.category_axis
    val_axis = chart.value_axis

    cat_axis.tick_labels.font.size = Pt(12)
    

    val_axis.minimum_scale = 0.0
    val_axis.maximum_scale = 100.0
    val_axis.major_unit = 20.0

    plot.has_data_labels = True
    labels = plot.data_labels
    labels.font.size = Pt(10)
    labels.number_format = '0"%"'
    _set_axis_titles(chart, "Difficulty", "Percentage")



def _format_perf_summary(chart, plot):

    plot.gap_width = 180

    cat_axis = chart.category_axis
    val_axis = chart.value_axis

    cat_axis.tick_labels.font.size = Pt(12)

    val_axis.minimum_scale = 0.0
    val_axis.maximum_scale = 100.0
    val_axis.major_unit = 20.0

    plot.has_data_labels = True
    labels = plot.data_labels
    labels.font.size = Pt(11)
    labels.number_format = '0"%"'
    _set_axis_titles(chart, "Category", "Average (%)")



def _format_subwise(chart, plot):
    print("SUBWISE FORMAT APPLIED")
    plot.gap_width = 180

    cat_axis = chart.category_axis
    val_axis = chart.value_axis

    # CONTRACT: Subwise charts are always 0-100%
    val_axis.minimum_scale = 0.0
    val_axis.maximum_scale = 100.0
    val_axis.major_unit = 20.0

    plot.has_data_labels = True
    labels = plot.data_labels
    labels.font.size = Pt(11)
    labels.number_format = '0"%"'
    _set_axis_titles(chart, "Subject", "Average (%)")



def _format_generic(chart, plot):

    plot.gap_width = 200

    plot.has_data_labels = True
    labels = plot.data_labels
    labels.font.size = Pt(9)
    _set_axis_titles(chart, "Category", "Value")



# --------------------------------------------------
# HELPERS
# --------------------------------------------------

def _apply_bar_colors(chart):

    for series in chart.series:
        for point in series.points:
            fill = point.format.fill
            fill.solid()
            fill.fore_color.rgb = DEFAULT_BAR_COLOR


def _set_axis_titles(chart, x_title=None, y_title=None):

    if x_title and chart.category_axis:
        axis = chart.category_axis
        axis.has_title = True
        tf = axis.axis_title.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = x_title
        p.font.size = Pt(11)
        p.font.bold = False
        p.font.color.rgb = RGBColor(50, 50, 50)

    if y_title and chart.value_axis:
        axis = chart.value_axis
        axis.has_title = True
        tf = axis.axis_title.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = y_title
        p.font.size = Pt(11)
        p.font.bold = False
        p.font.color.rgb = RGBColor(50, 50, 50)


def _render_debug_overlay(ppt_slide, x, y, w, h, element, dataset_type):
    """
    Renders a semi-transparent text box inside the chart area with debug info.
    """
    padding = px_to_emu(10)
    
    box = ppt_slide.shapes.add_textbox(
        x + padding, 
        y + padding, 
        px_to_emu(200), 
        px_to_emu(60)
    )
    
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 204)
    fill.transparency = 0.2
    
    line = box.line
    line.color.rgb = RGBColor(255, 0, 0)
    line.width = Pt(1)

    tf = box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    dataset_name = element.get("datasetName", "Unknown")
    chart_type = element.get("chartType", "Unknown")
    
    text = (
        f"DEBUG INFO:\n"
        f"Dataset: {dataset_name}\n"
        f"Type: {dataset_type} ({chart_type})"
    )
    p.text = text

